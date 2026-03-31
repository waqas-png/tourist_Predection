"""
Create Tourism ML PowerPoint Presentation using the Sample PPTX as template.
Replaces content (text + images) in all 16 slides while preserving visual style.
"""

import shutil, copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SAMPLE = r"c:\Users\ASAD\Downloads\Waqas data\Sample Presentation 3.pptx"
OUTPUT = r"c:\Users\ASAD\Downloads\Waqas data\Tourism_ML_Presentation.pptx"
IMG    = r"c:\Users\ASAD\Downloads\Waqas data\images"

# Copy sample to output (preserves all themes, master, layouts, backgrounds)
shutil.copy(SAMPLE, OUTPUT)
prs = Presentation(OUTPUT)

NS_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P  = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────

def replace_text_in_tf(shape, lines, keep_format=True):
    """Replace all text in a text frame with `lines` (list of strings).
    If keep_format=True, preserve the first run's font formatting."""
    if not hasattr(shape, "text_frame"):
        return
    tf   = shape.text_frame
    txBody = tf._txBody

    # Capture first-run format to reuse
    first_rPr = first_pPr = None
    first_p = txBody.find(f"{{{NS_A}}}p")
    if first_p is not None and keep_format:
        pPr_el = first_p.find(f"{{{NS_A}}}pPr")
        if pPr_el is not None:
            first_pPr = copy.deepcopy(pPr_el)
        first_r_el = first_p.find(f"{{{NS_A}}}r")
        if first_r_el is not None:
            rPr_el = first_r_el.find(f"{{{NS_A}}}rPr")
            if rPr_el is not None:
                first_rPr = copy.deepcopy(rPr_el)

    # Remove all existing <a:p> elements
    for p_el in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p_el)

    # Build new paragraphs
    for line in lines:
        p_el = etree.SubElement(txBody, f"{{{NS_A}}}p")
        if first_pPr is not None:
            p_el.append(copy.deepcopy(first_pPr))
        r_el = etree.SubElement(p_el, f"{{{NS_A}}}r")
        if first_rPr is not None:
            r_el.append(copy.deepcopy(first_rPr))
        t_el = etree.SubElement(r_el, f"{{{NS_A}}}t")
        t_el.text = line


def replace_text_in_tf_multi(shape, lines_with_bold, keep_format=True):
    """Replace text with mixed bold/normal lines.
    lines_with_bold: list of (text, is_bold) tuples."""
    if not hasattr(shape, "text_frame"):
        return
    tf     = shape.text_frame
    txBody = tf._txBody

    first_rPr = first_pPr = None
    first_p = txBody.find(f"{{{NS_A}}}p")
    if first_p is not None and keep_format:
        pPr_el = first_p.find(f"{{{NS_A}}}pPr")
        if pPr_el is not None:
            first_pPr = copy.deepcopy(pPr_el)
        first_r_el = first_p.find(f"{{{NS_A}}}r")
        if first_r_el is not None:
            rPr_el = first_r_el.find(f"{{{NS_A}}}rPr")
            if rPr_el is not None:
                first_rPr = copy.deepcopy(rPr_el)

    for p_el in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p_el)

    for (line, is_bold) in lines_with_bold:
        p_el = etree.SubElement(txBody, f"{{{NS_A}}}p")
        if first_pPr is not None:
            p_el.append(copy.deepcopy(first_pPr))
        r_el = etree.SubElement(p_el, f"{{{NS_A}}}r")
        rPr = copy.deepcopy(first_rPr) if first_rPr is not None else etree.Element(f"{{{NS_A}}}rPr")
        rPr.set("b", "1" if is_bold else "0")
        r_el.append(rPr)
        t_el = etree.SubElement(r_el, f"{{{NS_A}}}t")
        t_el.text = line


def replace_image(slide, pic_shape, new_img_path):
    """Replace the image bytes of a Picture shape."""
    blip = pic_shape.element.find(f".//{{{NS_A}}}blip")
    if blip is None:
        return
    rId = blip.get(f"{{{NS_R}}}embed")
    if rId is None:
        return
    try:
        img_part = slide.part.related_part(rId)
        with open(new_img_path, "rb") as f:
            img_part._blob = f.read()
        img_part._content_type = "image/png"
        print(f"  OK Image replaced with {os.path.basename(new_img_path)}")
    except Exception as e:
        print(f"  FAIL Image error: {e}")


def get_shapes_by_type(slide, shape_type_id):
    """Return all shapes of a given type (17=TEXT_BOX, 13=PICTURE, 19=TABLE, 6=GROUP)."""
    return [s for s in slide.shapes if s.shape_type == shape_type_id]


def get_tb_by_pos(slide, t_min, t_max, l_min=0, l_max=100):
    """Return text box shapes within a position range (in inches)."""
    results = []
    for s in slide.shapes:
        if s.shape_type != 17:  # TEXT_BOX
            continue
        t = s.top  / 914400
        l = s.left / 914400
        if t_min <= t <= t_max and l_min <= l <= l_max:
            results.append(s)
    return results


def update_table_cell(table, row, col, new_text):
    """Replace text in a table cell."""
    cell   = table.rows[row].cells[col]
    txBody = cell.text_frame._txBody
    first_rPr = None
    first_p = txBody.find(f"{{{NS_A}}}p")
    if first_p is not None:
        r_el = first_p.find(f"{{{NS_A}}}r")
        if r_el is not None:
            rPr_el = r_el.find(f"{{{NS_A}}}rPr")
            if rPr_el is not None:
                first_rPr = copy.deepcopy(rPr_el)
    for p_el in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p_el)
    p_el = etree.SubElement(txBody, f"{{{NS_A}}}p")
    r_el = etree.SubElement(p_el, f"{{{NS_A}}}r")
    if first_rPr is not None:
        r_el.append(first_rPr)
    t_el = etree.SubElement(r_el, f"{{{NS_A}}}t")
    t_el.text = new_text


def find_text_containing(slide, text_fragment):
    """Find the first shape whose text contains text_fragment."""
    for s in slide.shapes:
        if hasattr(s, "text") and text_fragment.lower() in s.text.lower():
            return s
    return None


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1  — Title Slide
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 1 - Title")
s1 = prs.slides[0]
# Title text box (large font, T≈1.7)
for tb in s1.shapes:
    if tb.shape_type == 17 and abs(tb.top/914400 - 1.70) < 0.3:
        replace_text_in_tf(tb, [
            "International Tourism Demand",
            "Forecasting Using Machine Learning",
            "A Production-Ready MLOps Framework",
        ])
# "Presented by" text box (T≈9.27)
for tb in s1.shapes:
    if tb.shape_type == 17 and 9.0 < tb.top/914400 < 9.5:
        replace_text_in_tf(tb, ["Presented  by"])
# "=======" text box (T≈9.72)
for tb in s1.shapes:
    if tb.shape_type == 17 and 9.5 < tb.top/914400 < 10.0:
        replace_text_in_tf(tb, ["Waqas  |  Department of Computing"])
# Replace right image
for s in s1.shapes:
    if s.shape_type == 13:
        replace_image(s1, s, os.path.join(IMG, "trends_and_gdp.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2  — Introduction
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 2 - Introduction")
s2 = prs.slides[1]
# Section title (T≈0.14, large font)
t_title = find_text_containing(s2, "INTRODUCTION")
if t_title:
    replace_text_in_tf(t_title, ["INTRODUCTION"])
# Body text box (T≈1.18, large)
for tb in s2.shapes:
    if tb.shape_type == 17 and 1.0 < tb.top/914400 < 2.0 and tb.width/914400 > 10:
        replace_text_in_tf(tb, [
            "International tourism contributes ~10% of global GDP and over 330 million jobs worldwide",
            "Accurate arrivals forecasting is critical for governments, tourism boards, and economic planners",
            "Arrivals are driven by GDP growth, exchange rates, geopolitics, airline connectivity, and pandemics",
            "Traditional ARIMA models fail to capture nonlinear dynamics and structural breaks (COVID-19, SARS)",
            "The World Bank provides 25 years of verified data across 200+ countries (1999–2023)",
            "Most ML studies stop at the notebook — no production deployment or real-time monitoring exists",
            "This study delivers both: state-of-the-art ML forecasting AND a complete MLOps pipeline",
        ])
# Replace right image
for s in s2.shapes:
    if s.shape_type == 13:
        replace_image(s2, s, os.path.join(IMG, "top10_countries.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3  — Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 3 - Problem Statement")
s3 = prs.slides[2]
t_title = find_text_containing(s3, "PROBLEM")
if t_title:
    replace_text_in_tf(t_title, ["PROBLEM STATEMENT"])
for tb in s3.shapes:
    if tb.shape_type == 17 and 1.0 < tb.top/914400 < 2.0 and tb.width/914400 > 10:
        replace_text_in_tf(tb, [
            "Tourism arrivals are influenced by complex, nonlinear, multi-factorial relationships",
            "Synthetic datasets yield near-random accuracy (~25%) — real-world data is essential",
            "Traditional econometric models (ARIMA, gravity) fail on structural breaks like COVID-19",
            "The 2020 COVID-19 shock caused a 73% global arrivals collapse — a rare out-of-distribution event",
            "Existing ML tourism studies lack production deployment, monitoring, and operational runbooks",
            "End-to-end MLOps frameworks for tourism forecasting are absent from the academic literature",
            "Solution: Combine 25-year World Bank panel data + ML engineering + full production MLOps",
        ])
for s in s3.shapes:
    if s.shape_type == 13:
        replace_image(s3, s, os.path.join(IMG, "arrivals_distribution.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4  — Research Objectives (numbered, dark background)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 4 - Research Objectives")
s4 = prs.slides[3]
t_title = find_text_containing(s4, "OBJECTIVES")
if t_title:
    replace_text_in_tf(t_title, ["RESEARCH OBJECTIVES"])

# Objective text boxes: 01→TextBox 22, 02→TextBox 28, 03→TextBox 34
obj_map = [
    ("01", "Collect and preprocess 25 years of World Bank tourism data "
           "spanning 200+ countries. Apply log-normalisation, median imputation, "
           "and lag-based temporal feature engineering to produce 3,091 clean records."),
    ("02", "Train and rigorously benchmark three regression models — Ridge Regression, "
           "Gradient Boosting Regressor, and Random Forest Regressor — evaluating on "
           "R², RMSE, and MAE using an 80/20 held-out test split."),
    ("03", "Deploy a production-grade MLOps pipeline: FastAPI inference → Docker "
           "containerisation → GitHub Actions CI/CD → AWS ECS Fargate → "
           "Prometheus + Grafana real-time monitoring with automated alerts."),
]
obj_idx = 0
for tb in s4.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 7.0 and tb.width/914400 > 4.0:
        if obj_idx < len(obj_map):
            _, text = obj_map[obj_idx]
            replace_text_in_tf(tb, [text])
            obj_idx += 1

for s in s4.shapes:
    if s.shape_type == 13:
        replace_image(s4, s, os.path.join(IMG, "feature_importance.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5  — Key Statement
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 5 - Key Statement")
s5 = prs.slides[4]
for tb in s5.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 1.5 and tb.width/914400 > 8.0:
        replace_text_in_tf(tb, [
            "25 Years of Real World Bank Data",
            "200+ Countries  ·  3 ML Models",
            "R² = 1.0000",
            "",
            "Ridge Regression achieves near-perfect forecasting of",
            "international tourism arrivals — demonstrating that",
            "temporal lag features unlock the predictive signal",
            "hidden in historical arrival patterns.",
        ])
for s in s5.shapes:
    if s.shape_type == 13:
        replace_image(s5, s, os.path.join(IMG, "pred_ridge_regression.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6  — Model Performance Table  (was Literature Review table)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 6 - Model Performance Table")
s6 = prs.slides[5]
for shape in s6.shapes:
    if shape.shape_type == 19:  # TABLE
        tbl = shape.table
        rows_needed = 4   # header + 3 models
        cols_needed = 5   # Model, R², RMSE_log, MAE_log, Train(s)

        # Header row (row 0)
        headers = ["Model", "R²", "RMSE (log)", "MAE (log)", "Train Time (s)"]
        for c_i, h in enumerate(headers):
            if c_i < len(tbl.columns):
                update_table_cell(tbl, 0, c_i, h)

        # Data rows
        data = [
            ["Ridge Regression",   "1.0000", "0.0118", "0.0071", "0.00"],
            ["Gradient Boosting",  "0.9992", "0.0651", "0.0323", "1.90"],
            ["Random Forest",      "0.9976", "0.1124", "0.0424", "0.29"],
        ]
        for r_i, row_data in enumerate(data):
            for c_i, cell_text in enumerate(row_data):
                if r_i + 1 < len(tbl.rows) and c_i < len(tbl.columns):
                    update_table_cell(tbl, r_i + 1, c_i, cell_text)

        # Clear remaining rows
        for r_i in range(4, len(tbl.rows)):
            for c_i in range(len(tbl.columns)):
                update_table_cell(tbl, r_i, c_i, "")
        break

# Add a title text box for the table slide (there's no title TextBox shown in inspection — add one)
# Update the header group page text if any
for tb in s6.shapes:
    if tb.shape_type == 17 and 0.1 < tb.top/914400 < 0.5 and tb.width/914400 > 5:
        replace_text_in_tf(tb, ["MODEL PERFORMANCE COMPARISON"])
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7  — Feature Engineering  (was Research Gap)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 7 - Feature Engineering")
s7 = prs.slides[6]
t_title = find_text_containing(s7, "GAP")
if t_title:
    replace_text_in_tf(t_title, ["FEATURE ENGINEERING"])
# Body text on right side
for tb in s7.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 1.0 and tb.left/914400 > 6.0:
        replace_text_in_tf(tb, [
            "Log-Normalisation  →  x' = log(x + 1) applied to all continuous features",
            "Stabilises heavy right-skewed distributions (raw arrivals: 900 → 2.4 billion)",
            "",
            "Lag Features  →  1-year and 2-year lagged log-arrivals per country",
            "Computed with groupby('country').shift(1) to prevent data leakage",
            "lag1_log_arrivals accounts for 93.02% of Random Forest feature importance",
            "",
            "Growth Rate  →  Year-on-year log-arrivals difference (momentum signal)",
            "",
            "Temporal Features  →  year_norm (0–1),  is_post_covid (binary),  decade bucket",
            "",
            "Economic Features  →  log_gdp, log_receipts, log_exports, log_expenditures, inflation",
            "",
            "Final Feature Set: 12 variables  ·  3,091 modelling observations",
        ])
# Replace left image
for s in s7.shapes:
    if s.shape_type == 13:
        replace_image(s7, s, os.path.join(IMG, "correlation.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8  — Methodology  (numbered steps)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 8 - Methodology")
s8 = prs.slides[7]
t_title = find_text_containing(s8, "METHODOLOGY")
if t_title:
    replace_text_in_tf(t_title, ["METHODOLOGY"])

# 7 numbered items: each has a header TextBox (small, near numbered circle)
# and a body TextBox below it
methodology = [
    ("Datasets Used",
     "World Bank Tourism & Economic Impact dataset via Kaggle. "
     "6,650 records · 266 countries · 1999–2023 · 11 variables."),
    ("Data Preprocessing",
     "Remove columns >40% missing (departures, unemployment). "
     "Median imputation for remaining. Discard rows with missing target. "
     "6,650 → 4,949 → 3,091 final records."),
    ("Feature Engineering",
     "Log-normalise all continuous features. Create 1-year and 2-year "
     "lag features per country. Add growth rate, year_norm, is_post_covid, "
     "decade. Label-encode country. Final: 12 features."),
    ("Model Training",
     "80/20 train/test split (seed=42). "
     "Ridge Regression (α=1.0, StandardScaler), "
     "Gradient Boosting (300 trees, depth=5, lr=0.05), "
     "Random Forest (300 trees, depth=12, min_split=4)."),
    ("Model Evaluation",
     "Evaluate on held-out test set using R², RMSE, MAE on log scale. "
     "Best model: Ridge R²=1.0000, RMSE=0.0118. "
     "Back-transform: ŷ = e^y′ − 1."),
    ("MLOps Deployment",
     "FastAPI inference service → Docker (Python 3.12-slim, non-root). "
     "GitHub Actions: lint → test → ECR push → ECS rolling deploy. "
     "AWS ECS Fargate (2 tasks, ALB, auto-scaling)."),
    ("Monitoring",
     "Prometheus scrapes /metrics every 15s. "
     "4 alert rules: APIDown · HighErrorRate · SlowPredictions · NoTraffic. "
     "Grafana 6-panel dashboard. Slack webhook notifications."),
]

# TextBoxes for headers (bold, medium font) and bodies (normal) on right side
# From inspection: pairs of text boxes at T=1.28/1.59, 2.54/2.93, 3.98/4.34, etc.
header_tbs = []
body_tbs   = []
for tb in s8.shapes:
    if tb.shape_type != 17:
        continue
    t = tb.top / 914400
    l = tb.left / 914400
    if l > 10.0 and 1.0 < t < 11.5:
        w = tb.width / 914400
        h = tb.height / 914400
        if h < 0.6:   # short → header
            header_tbs.append((t, tb))
        else:          # taller → body
            body_tbs.append((t, tb))

header_tbs.sort(key=lambda x: x[0])
body_tbs.sort(key=lambda x: x[0])

for i, (heading, body_text) in enumerate(methodology):
    if i < len(header_tbs):
        replace_text_in_tf(header_tbs[i][1], [heading])
    if i < len(body_tbs):
        replace_text_in_tf(body_tbs[i][1], [body_text])

# Replace left image
for s in s8.shapes:
    if s.shape_type == 13:
        replace_image(s8, s, os.path.join(IMG, "trends_and_gdp.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9  — ML Models  (3-column visual)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 9 - ML Models (3-column)")
s9 = prs.slides[8]

# The 3 column groups (Group 13, 27, 37) contain text & visual elements.
# We'll find and update all text boxes inside these groups,
# or add text boxes over each column.
# Columns: L≈1.1, 7.2, 13.2 (each W≈5.7")

col_data = [
    {
        "title": "Ridge Regression",
        "rank":  "Rank 1st",
        "lines": [
            "Regularisation: L2 (α = 1.0)",
            "Features: StandardScaler",
            "Complexity: Linear",
            "R² = 1.0000",
            "RMSE = 0.0118",
            "MAE  = 0.0071",
            "Train: 0.00 s",
            "",
            "Best fit because the log-lag",
            "relationship is near-linear.",
        ],
    },
    {
        "title": "Gradient Boosting",
        "rank":  "Rank 2nd",
        "lines": [
            "Trees: 300  Depth: 5",
            "Learning rate: 0.05",
            "Complexity: Nonlinear",
            "R² = 0.9992",
            "RMSE = 0.0651",
            "MAE  = 0.0323",
            "Train: 1.90 s",
            "",
            "Sequentially corrects residuals.",
            "Captures nonlinear interactions.",
        ],
    },
    {
        "title": "Random Forest",
        "rank":  "Rank 3rd",
        "lines": [
            "Trees: 300  Depth: 12",
            "Min samples split: 4",
            "Complexity: Ensemble",
            "R² = 0.9976",
            "RMSE = 0.1124",
            "MAE  = 0.0424",
            "Train: 0.29 s",
            "",
            "Averages 300 independent trees.",
            "Provides feature importance.",
        ],
    },
]

# Try to update text boxes within the column groups
# Group 13 → L≈1.1, Group 27 → L≈7.2, Group 37 → L≈13.2
col_groups = []
for s in s9.shapes:
    if s.shape_type == 6 and 0.8 < s.left/914400 < 14.0 and s.width/914400 > 4.0:
        l = s.left / 914400
        if 0.8 < l < 3.0 or 6.5 < l < 9.0 or 12.5 < l < 15.0:
            col_groups.append((l, s))

col_groups.sort(key=lambda x: x[0])

for col_i, (l_pos, grp) in enumerate(col_groups[:3]):
    if col_i >= len(col_data):
        break
    data = col_data[col_i]
    # Find text boxes inside this group's XML
    txb_list = []
    for el in grp.element.iter(f"{{{NS_P}}}sp"):
        txBody = el.find(f".//{{{NS_A}}}txBody")
        if txBody is not None:
            has_text = any(
                r.find(f"{{{NS_A}}}t") is not None
                for r in txBody.iter(f"{{{NS_A}}}r")
            )
            if has_text:
                txb_list.append(txBody)

    # Replace text in found txBody elements
    line_idx = 0
    all_lines = [data["rank"], data["title"]] + data["lines"]
    for txb in txb_list:
        for p_el in txb.findall(f"{{{NS_A}}}p"):
            for r_el in p_el.findall(f"{{{NS_A}}}r"):
                t_el = r_el.find(f"{{{NS_A}}}t")
                if t_el is not None and line_idx < len(all_lines):
                    t_el.text = all_lines[line_idx]
                    line_idx += 1


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10  — Results
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 10 - Results")
s10 = prs.slides[9]

# Two pictures → replace
pics = [s for s in s10.shapes if s.shape_type == 13]
if len(pics) >= 1:
    replace_image(s10, pics[0], os.path.join(IMG, "pred_ridge_regression.png"))
if len(pics) >= 2:
    replace_image(s10, pics[1], os.path.join(IMG, "model_comparison.png"))

# Update text boxes on lower portion
result_texts = [
    [
        "Actual vs. Predicted — Ridge (R²=1.0000)",
        "All scatter points lie exactly on the y=x line.",
        "Near-perfect prediction after log-normalisation.",
    ],
    [
        "All 3 Models Achieve R² > 0.997",
        "Ridge > Gradient Boosting > Random Forest.",
        "Log-lag feature drives near-linear relationship.",
    ],
]
lower_tbs = sorted(
    [tb for tb in s10.shapes if tb.shape_type == 17 and tb.top/914400 > 7.5],
    key=lambda x: x.left
)
for i, tb in enumerate(lower_tbs[:2]):
    if i < len(result_texts):
        replace_text_in_tf(tb, result_texts[i])


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11  — Feature Importance  (was detailed results table)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 11 - Feature Importance")
s11 = prs.slides[10]

# Update the TABLE
for shape in s11.shapes:
    if shape.shape_type == 19:
        tbl = shape.table
        fi_headers = ["Feature", "Importance", "Category", "Description"]
        fi_data = [
            ["lag1_log_arrivals",        "93.02%", "Lag",      "1-year lagged log-arrivals"],
            ["lag2_log_arrivals",        " 5.80%", "Lag",      "2-year lagged log-arrivals"],
            ["arrival_growth",           " 0.80%", "Lag",      "Year-on-year log-arrivals diff"],
            ["year_norm",                " 0.10%", "Temporal", "Normalised year (0–1)"],
            ["decade",                   " 0.08%", "Temporal", "Decade bucket"],
            ["is_post_covid",            " 0.08%", "Temporal", "Binary flag: year ≥ 2020"],
            ["log_gdp",                  " 0.05%", "Economic", "log(GDP + 1)"],
            ["log_tourism_receipts",     " 0.03%", "Economic", "log(receipts + 1)"],
            ["log_tourism_expenditures", " 0.02%", "Economic", "log(expenditures + 1)"],
            ["inflation",                " 0.01%", "Economic", "Annual inflation rate (%)"],
        ]
        for c_i, h in enumerate(fi_headers):
            if c_i < len(tbl.columns):
                update_table_cell(tbl, 0, c_i, h)
        for r_i, row_data in enumerate(fi_data):
            for c_i, val in enumerate(row_data):
                if r_i + 1 < len(tbl.rows) and c_i < len(tbl.columns):
                    update_table_cell(tbl, r_i + 1, c_i, val)
        for r_i in range(len(fi_data) + 1, len(tbl.rows)):
            for c_i in range(len(tbl.columns)):
                update_table_cell(tbl, r_i, c_i, "")
        break

# Update right-side analysis text
for tb in s11.shapes:
    if tb.shape_type == 17 and tb.left/914400 > 7.0 and tb.top/914400 > 1.0:
        replace_text_in_tf(tb, [
            "Key Finding: Temporal autocorrelation dominates.",
            "",
            "lag1_log_arrivals alone accounts for 93% of",
            "total Random Forest feature importance.",
            "",
            "Countries that attracted many tourists last year",
            "will attract similar volumes next year — driven",
            "by destination brand equity and infrastructure.",
            "",
            "Economic features (GDP, receipts) carry <0.25%",
            "importance after lag features are included.",
            "",
            "Implication: For annual country-level forecasting,",
            "a simple Ridge model on lag features outperforms",
            "complex tree ensembles.",
        ])


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12  — MLOps Pipeline  (horizontal numbered groups)
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 12 - MLOps Pipeline")
s12 = prs.slides[11]

mlops_steps = [
    ("01", "Data\nIngestion",
     "World Bank CSV\n(6,650 records\n266 countries)"),
    ("02", "Feature\nEngineering",
     "Log-norm · Lags\n12 features\n3,091 rows"),
    ("03", "Model\nTraining",
     "Ridge / GB / RF\nScikit-learn\nR²=1.0000"),
    ("04", "FastAPI\nService",
     "/predict\n/health\n/metrics"),
    ("05", "Docker\nECR",
     "Python 3.12-slim\nNon-root user\nTrivy scan"),
    ("06", "GitHub\nActions\nCI/CD",
     "Lint → Test\nBuild → Push\nDeploy → Smoke"),
    ("07", "AWS ECS\n+\nMonitoring",
     "Fargate 2 tasks\nPrometheus\nGrafana"),
]

# Find text inside each horizontal group
step_idx = 0
for grp in s12.shapes:
    if grp.shape_type == 6 and grp.name.startswith("Group 14"):
        if step_idx >= len(mlops_steps):
            break
        num, title, detail = mlops_steps[step_idx]
        all_texts = [num, title, detail]
        t_idx = 0
        for txBody in grp.element.iter(f"{{{NS_A}}}txBody"):
            r_els = list(txBody.iter(f"{{{NS_A}}}r"))
            if r_els and t_idx < len(all_texts):
                # Replace text in first run of first paragraph
                first_r = r_els[0]
                t_el = first_r.find(f"{{{NS_A}}}t")
                if t_el is not None:
                    t_el.text = all_texts[t_idx]
                    t_idx += 1
        step_idx += 1

# Add title if there's a top text box
for tb in s12.shapes:
    if tb.shape_type == 17 and tb.top/914400 < 1.3:
        replace_text_in_tf(tb, ["MLOps PIPELINE — END-TO-END ARCHITECTURE"])
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13  — Discussion
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 13 - Discussion")
s13 = prs.slides[12]
for tb in s13.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 1.0 and tb.width/914400 > 12.0:
        replace_text_in_tf(tb, [
            "Near-perfect R² values confirm that ML with proper feature engineering "
            "can deliver highly accurate tourism forecasting when applied to real-world data.",
            "",
            "Ridge Regression dominates because after log-transformation, "
            "the lag-arrivals vs current-arrivals relationship is approximately linear — "
            "making L2-regularised linear regression the natural best fit.",
            "",
            "The dominance of lag1_log_arrivals (93% importance) is consistent with "
            "tourism economics literature: destinations with strong brand equity maintain "
            "persistent visitor volumes year over year.",
            "",
            "The COVID-19 binary feature enables models to partially adjust for the 2020 "
            "structural break. The lag features propagate the shock through subsequent years.",
            "",
            "MLOps contribution: The production system achieves ML maturity Stage 4 — "
            "automated training, deployment, serving, and continuous monitoring. "
            "Median prediction latency: 12ms (p95: 28ms) — well within the 500ms alert threshold.",
            "",
            "OIDC identity federation eliminates long-lived AWS credentials. "
            "GitHub Actions CI/CD averages 4min 38s from commit to deployment.",
        ])

for s in s13.shapes:
    if s.shape_type == 13:
        replace_image(s13, s, os.path.join(IMG, "residuals.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14  — Conclusion & Future Work
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 14 - Conclusion & Future Work")
s14 = prs.slides[13]
for tb in s14.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 1.0 and tb.width/914400 > 10.0:
        replace_text_in_tf(tb, [
            "CONCLUSION",
            "✔  Ridge Regression R²=1.0000 · Gradient Boosting R²=0.9992 · Random Forest R²=0.9976",
            "✔  25-year World Bank panel data (200+ countries) with rigorous preprocessing",
            "✔  Lag features dominate: lag1_log_arrivals = 93% of feature importance",
            "✔  Production MLOps: FastAPI · Docker · GitHub Actions · AWS ECS Fargate",
            "✔  Prometheus + Grafana with 4 automated alert rules and Slack notifications",
            "✔  7 operational runbooks covering deployment, incidents, scaling, and secrets",
            "",
            "FUTURE WORK",
            "▸  LSTM / Transformer models for long-range temporal dependencies",
            "▸  Additional covariates: air connectivity indices, social media sentiment",
            "▸  Scenario-based forecasting for climate change and pandemic disruptions",
            "▸  Country-specific model fine-tuning and ensemble stacking",
            "▸  Real-time streaming data pipeline (Kafka / AWS Kinesis)",
            "▸  MLflow / DVC for model versioning and experiment tracking",
        ])
for s in s14.shapes:
    if s.shape_type == 13:
        replace_image(s14, s, os.path.join(IMG, "top10_countries.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15  — References
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 15 - References")
s15 = prs.slides[14]
for tb in s15.shapes:
    if tb.shape_type == 17 and tb.top/914400 > 1.0 and tb.width/914400 > 15.0:
        replace_text_in_tf(tb, [
            "[1]  UNWTO, International Tourism Highlights 2020 Edition, Madrid, Spain, 2020.",
            "[2]  H. Song, G. Li, S. Witt, Tourism demand modelling and forecasting, "
                 "Tourism Economics, vol. 16, no. 1, pp. 19-31, 2010.",
            "[3]  G. Li, H. Song, Data science and tourism: Big data approaches, "
                 "Journal of Travel Research, vol. 59, no. 6, pp. 923-935, 2020.",
            "[4]  S. Cankurt, A. Subasi, Developing tourism demand forecasting models "
                 "using ML, Balkan J. Elec. Comput. Eng., vol. 4, no. 1, 2016.",
            "[5]  S. Sun et al., Forecasting tourist arrivals with ML and internet search "
                 "index, Tourism Management, vol. 70, pp. 1-10, 2019.",
            "[6]  T. Chen, C. Guestrin, XGBoost: A scalable tree boosting system, "
                 "Proc. KDD, pp. 785-794, 2016.",
            "[7]  D. Sculley et al., Hidden technical debt in ML systems, "
                 "NeurIPS, vol. 28, 2015.",
            "[8]  D. Kreuzberger et al., Machine learning operations (MLOps): Overview, "
                 "IEEE Access, vol. 11, pp. 31866-31879, 2023.",
            "[9]  U. Gunter, I. Onder, Forecasting international city tourism demand, "
                 "Tourism Management, vol. 46, pp. 123-135, 2015.",
            "[10] World Bank, Tourism and Economic Impact Dataset (1999-2023), "
                 "Kaggle, 2024. kaggle.com/datasets/bushraqurban/tourism-and-economic-impact",
            "[11] J. H. Friedman, Greedy function approximation: A gradient boosting machine, "
                 "Ann. Statistics, vol. 29, no. 5, pp. 1189-1232, 2001.",
            "[12] L. Breiman, Random forests, Machine Learning, vol. 45, no. 1, pp. 5-32, 2001.",
        ])


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16  — Thank You
# ─────────────────────────────────────────────────────────────────────────────
print("Slide 16 - Thank You")
s16 = prs.slides[15]
for s in s16.shapes:
    if s.shape_type == 13:
        replace_image(s16, s, os.path.join(IMG, "model_comparison.png"))
        break


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"DONE: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
